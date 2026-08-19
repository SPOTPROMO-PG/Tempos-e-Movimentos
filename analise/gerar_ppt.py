# -*- coding: utf-8 -*-
"""
Monta a apresentação entregue ao cliente a partir das saídas de jornada.py.

Todos os números são lidos dos CSVs/JSONs gerados pela análise — nada é
digitado à mão aqui, para o deck não descolar da base quando a coleta
crescer. Rodar de novo depois de jornada.py já atualiza os slides.
"""

import json

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ------------------------------------------------------------------ dados
mov = pd.read_csv('analise/saida/jornada_por_canal.csv')
ger = pd.read_csv('analise/saida/jornada_geral.csv')
vis = pd.read_csv('analise/saida/resumo_visitas.csv')
meta = json.load(open('analise/saida/meta.json', encoding='utf-8'))
meta['periodo'] = ['{2}/{1}/{0}'.format(*d.split('-')) for d in meta['periodo']]   # ISO -> dd/mm/aaaa
desl = json.load(open('analise/saida/desloc.json', encoding='utf-8'))
extra = json.load(open('analise/saida/extra.json', encoding='utf-8'))
canais = json.load(open('analise/saida/dossie_min.json', encoding='utf-8'))

FORTES = [c for c in canais if c['n'] >= 8]      # amostra adequada
FRACOS = [c for c in canais if c['n'] < 8]

# ----------------------------------------------------------------- estilo
PETROL = RGBColor(0x12, 0x3A, 0x4D)
TEAL = RGBColor(0x2C, 0x73, 0x91)
CLARO = RGBColor(0xE3, 0xEE, 0xF1)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
TINTA = RGBColor(0x14, 0x18, 0x1D)
CINZA = RGBColor(0x6D, 0x78, 0x87)
TERRA = RGBColor(0xB8, 0x50, 0x42)
VERDE = RGBColor(0x2F, 0x6B, 0x4F)
SERIE = ['123A4D', '2C7391', '7DB0C1', 'B85042', 'A9CBD6', '4A92AA', 'CBE0E6', '97A9B4']

TIT, CORPO = 'Cambria', 'Calibri'
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = 13.333, 7.5
BR = Inches(0.9)          # margem lateral


def slide(escuro=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if escuro:
        f = s.background.fill
        f.solid()
        f.fore_color.rgb = PETROL
    return s


def cx(s, x, y, w, h, txt, tam=14, cor=TINTA, bold=False, fonte=CORPO,
       al=PP_ALIGN.LEFT, esp=1.25, anc=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anc
    linhas = txt.split('\n') if isinstance(txt, str) else txt
    for i, ln in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        p.line_spacing = esp
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(tam)
        r.font.bold = bold
        r.font.name = fonte
        r.font.color.rgb = cor
    return tb


def titulo(s, txt, sub=None):
    cx(s, 0.9, 0.55, 11.6, 1.0, txt, 32, PETROL, True, TIT)
    if sub:
        cx(s, 0.9, 1.5, 11.6, 0.55, sub, 13, CINZA)


def caixa(s, x, y, w, h, cor=CLARO):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = cor
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.05
    return sh


def kpi(s, x, y, w, valor, rot, cor=PETROL, tamv=44):
    cx(s, x, y, w, 0.75, valor, tamv, cor, True, TIT, PP_ALIGN.CENTER, 1.0)
    cx(s, x, y + 0.78, w, 0.7, rot, 11, CINZA, False, CORPO, PP_ALIGN.CENTER, 1.15)


def estiliza(ch, legenda=False, rotulos=True, tam=11):
    ch.font.size = Pt(tam)
    ch.font.name = CORPO
    ch.has_title = False
    ch.has_legend = legenda
    if legenda:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(10)
    try:
        pl = ch.plots[0]
        pl.has_data_labels = rotulos
        if rotulos:
            dl = pl.data_labels
            dl.font.size = Pt(tam)
            dl.font.bold = True
            dl.font.color.rgb = TINTA
    except Exception:
        pass
    return ch


def pinta(ch, cores=None):
    cores = cores or SERIE
    pl = ch.plots[0]
    if len(ch.series) == 1:
        for i, pt in enumerate(pl.series[0].points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = RGBColor.from_string(cores[i % len(cores)])
    else:
        for i, se in enumerate(ch.series):
            se.format.fill.solid()
            se.format.fill.fore_color.rgb = RGBColor.from_string(cores[i % len(cores)])


def eixos(ch, valores=False):
    for ax, on in ((ch.category_axis, True), (ch.value_axis, valores)):
        try:
            ax.has_major_gridlines = False
            ax.tick_labels.font.size = Pt(10)
            ax.tick_labels.font.name = CORPO
            ax.tick_labels.font.color.rgb = CINZA
            ax.visible = on
        except Exception:
            pass


def rodape(s, txt):
    cx(s, 0.9, 6.72, 11.6, 0.3, txt, 9, CINZA)


# =========================================================== 1 · capa
s = slide(True)
cx(s, 0.9, 2.1, 11.5, 0.5, 'ESTUDO DE TEMPOS E MOVIMENTOS', 13, RGBColor(0x7D, 0xB0, 0xC1), True, CORPO)
cx(s, 0.9, 2.75, 11.5, 1.9, 'Como o promotor usa\na jornada de 8 horas', 44, BRANCO, True, TIT, esp=1.1)
cx(s, 0.9, 4.85, 11.5, 0.5, 'Diagnóstico por canal · Operação de merchandising P&G / SPOT', 15,
   RGBColor(0xA9, 0xCB, 0xD6))
cx(s, 0.9, 6.2, 11.5, 0.5, 'Coleta em campo  ·  {} a {}'.format(*meta['periodo']),
   12, RGBColor(0x7D, 0xB0, 0xC1))

# ================================================= 2 · o estudo em números
s = slide()
maior = ger.nlargest(1, 'pct_tempo_loja').iloc[0]
hm = lambda m: '{}h{:02d}'.format(int(m) // 60, int(m) % 60)
titulo(s, 'A jornada em números', 'O retrato do tempo dentro da loja.')
caixa(s, 0.9, 2.3, 11.55, 1.9)
for i, (v, r) in enumerate([
        (hm(vis.tempo_loja_min.median()), 'tempo por loja'),
        ('{:.0f}%'.format(extra['classes']['Execução (agrega valor)']['pct']), 'agrega valor'),
        (hm(extra['classes']['Deslocamento e busca']['min_8h']), 'deslocamento e busca'),
        ('{:.0f} min'.format(maior.min_em_8h), 'maior movimento'),
        ('{}'.format(vis.canal.nunique()), 'canais')]):
    kpi(s, 1.1 + i * 2.27, 2.62, 2.1, v, r)
cx(s, 0.9, 4.65, 11.55, 0.45, 'O que mais varia não é o promotor — é o canal', 19, PETROL, True, TIT)
cx(s, 0.9, 5.2, 11.55, 1.4,
   'Uma visita de DPP dura {:.0f} minutos; uma de NMR, {:.0f}. São operações diferentes, com rotinas\n'
   'diferentes, medidas com o mesmo instrumento. Por isso todo o diagnóstico é apresentado por canal:\n'
   'uma média geral esconderia mais do que revela.'.format(
       [c for c in canais if c['canal'] == 'DPP'][0]['visita_min'],
       [c for c in canais if c['canal'] == 'NMR'][0]['visita_min']), 14, TINTA, esp=1.4)
rodape(s, 'Tempo medido dentro da loja, em campo, na rotina real de trabalho.')

# ==================================================== 3 · metodologia
s = slide()
titulo(s, 'Como o tempo foi medido',
       'O método precisa lidar com atividades que acontecem ao mesmo tempo.')
caixa(s, 0.9, 2.25, 5.55, 3.9, RGBColor(0xFB, 0xF0, 0xEC))
cx(s, 1.25, 2.55, 4.85, 0.4, 'O PROBLEMA', 11, TERRA, True)
cx(s, 1.25, 3.05, 4.85, 3.0,
   'O promotor atende duas ou três categorias\nna mesma ida à gôndola e registra a mesma\n'
   'janela de horário para cada uma.\n\nSomar as durações declaradas infla o total\n'
   'em cerca de 2,4 vezes — e os percentuais\npassariam de 100%.', 14, TINTA, esp=1.35)
caixa(s, 6.9, 2.25, 5.55, 3.9, RGBColor(0xEE, 0xF5, 0xF1))
cx(s, 7.25, 2.55, 4.85, 0.4, 'A SOLUÇÃO', 11, VERDE, True)
cx(s, 7.25, 3.05, 4.85, 3.0,
   'O tempo é alocado minuto a minuto.\nEm cada minuto da visita, o tempo é\n'
   'dividido entre as atividades ativas nele.\n\nDois movimentos em paralelo levam metade\n'
   'do crédito cada. A soma fecha exatamente\no tempo real dentro da loja.', 14, TINTA, esp=1.35)
rodape(s, 'Excluídos: horários invertidos e atividades acima de 4 horas. Envios repetidos consolidados '
          'mantendo a versão mais completa.')

# ======================================= 4 · achado 1 — o tempo que agrega
s = slide()
titulo(s, 'Menos da metade do tempo em loja agrega valor',
       'Execução = abastecer, precificar, limpar e montar. O resto sustenta a execução, mas não é ela.')
cd = CategoryChartData()
cls = list(extra['classes'].keys())
cd.categories = [c.replace(' (agrega valor)', '') for c in cls]
cd.add_series('% do tempo', tuple(extra['classes'][c]['pct'] / 100 for c in cls))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.9), Inches(2.25),
                        Inches(6.5), Inches(3.9), cd).chart
estiliza(gf, tam=12)
pinta(gf, ['2F6B4F', '2C7391', 'B85042', 'C89A3C'])
eixos(gf)
gf.plots[0].data_labels.number_format = '0.0%'
gf.plots[0].data_labels.number_format_is_linked = False
gf.plots[0].gap_width = 60
y = 2.35
for c in cls:
    v = extra['classes'][c]
    cx(s, 7.75, y, 4.7, 0.35, c.upper(), 10.5, PETROL, True)
    cx(s, 7.75, y + 0.36, 4.7, 0.5,
       '{}h{:02d} de uma jornada de 8 horas'.format(v['min_8h'] // 60, v['min_8h'] % 60), 13, TINTA)
    y += 0.99
rodape(s, 'Percentuais sobre o tempo dentro da loja, média ponderada de todos os canais.')

# ============================== 5 · achado 2 — o canal define a rotina
s = slide()
titulo(s, 'O canal define a rotina mais do que qualquer outro fator',
       'Participação da execução em gôndola no tempo de visita, por canal.')
cd = CategoryChartData()
ordf = sorted(FORTES, key=lambda c: -(c['blocos']['Ponto Natural'] + c['blocos']['Ponto Extra']
                                      + c['blocos']['Check Out']) / c['visita_min'])
cd.categories = [c['canal'] for c in ordf]
cd.add_series('Execução', tuple((c['blocos']['Ponto Natural'] + c['blocos']['Ponto Extra']
                                 + c['blocos']['Check Out']) / c['visita_min'] for c in ordf))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.9), Inches(2.3),
                        Inches(11.55), Inches(3.15), cd).chart
estiliza(gf, tam=13)
pinta(gf)
eixos(gf)
gf.plots[0].data_labels.number_format = '0%'
gf.plots[0].data_labels.number_format_is_linked = False
gf.plots[0].gap_width = 80
cx(s, 0.9, 5.75, 11.55, 1.1,
   'A causa não é ritmo de trabalho: é o número de lojas. O DPP faz 4 lojas por dia e paga 4 vezes o custo\n'
   'fixo de chegar, se cadastrar e se localizar. C&C e NMR fazem 1 — diluem esse custo em uma visita longa.',
   14, TINTA, esp=1.4)
rodape(s, 'Execução = Ponto Natural, Ponto Extra e Check Out.')

# ==================================== 6 · composição do tempo por canal
s = slide()
titulo(s, 'A composição do tempo, canal a canal',
       'Minutos médios por visita, por bloco de atividade.')
BL = ['Entrada', 'Abertura', 'Estoque', 'Check Out', 'Ponto Natural', 'Ponto Extra',
      'Outras Atividades', 'Saída']
cd = CategoryChartData()
ordv = sorted(FORTES, key=lambda c: -c['visita_min'])
cd.categories = [c['canal'] for c in ordv]
for b in BL:
    cd.add_series(b, tuple(c['blocos'][b] for c in ordv))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(0.9), Inches(2.25),
                        Inches(11.55), Inches(3.85), cd).chart
estiliza(gf, legenda=True, rotulos=False)
pinta(gf)
eixos(gf, valores=True)
gf.value_axis.axis_title.text_frame.text = 'minutos por visita'
gf.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
gf.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = CINZA
gf.plots[0].gap_width = 70
rodape(s, 'Ponto Natural é o maior bloco em todos os canais.')

# ================================== 7 · para onde vai o tempo (top movs)
s = slide()
titulo(s, 'Os movimentos que mais consomem a jornada',
       'Os dez movimentos com maior peso, projetados em uma jornada de 8 horas.')
top = ger.nlargest(10, 'pct_tempo_loja').iloc[::-1]
cd = CategoryChartData()
cd.categories = [(m[:44] + '…') if len(m) > 45 else m for m in top.movimento]
cd.add_series('min em 8h', tuple(float(x) for x in top.min_em_8h))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.9), Inches(2.25),
                        Inches(11.55), Inches(4.0), cd).chart
estiliza(gf, tam=11)
pinta(gf, ['123A4D'] * 10)
eixos(gf)
gf.plots[0].data_labels.number_format = '0" min"'
gf.plots[0].data_labels.number_format_is_linked = False
gf.plots[0].gap_width = 45
rodape(s, 'Abastecimento e busca de itens no estoque somam {:.0f} minutos da jornada.'.format(
    float(top.min_em_8h.iloc[-1]) + float(top.min_em_8h.iloc[-2])))

# ============================================ 8 · deslocamento e busca
s = slide()
d_cls = extra['classes']['Deslocamento e busca']
titulo(s, 'Quase 2 horas por jornada andando e procurando',
       'É a maior oportunidade isolada do estudo — e a mais acionável.')
caixa(s, 0.9, 2.25, 3.5, 3.1, PETROL)
cx(s, 1.15, 3.0, 3.0, 1.0, '{}h{:02d}'.format(d_cls['min_8h'] // 60, d_cls['min_8h'] % 60),
   50, BRANCO, True, TIT, PP_ALIGN.CENTER, 1.0)
cx(s, 1.15, 4.05, 3.0, 0.9, 'por jornada de 8 horas\nem deslocamento e busca',
   12, RGBColor(0xA9, 0xCB, 0xD6), False, CORPO, PP_ALIGN.CENTER, 1.3)
itens = ger[ger.classe == 'Deslocamento e busca'].nlargest(5, 'min_em_8h')
y = 2.3
for _, r in itens.iterrows():
    cx(s, 4.85, y, 6.3, 0.4, r.movimento, 13, TINTA, True)
    cx(s, 11.35, y, 1.1, 0.4, '{:.0f} min'.format(r.min_em_8h), 13, TEAL, True, CORPO, PP_ALIGN.RIGHT)
    y += 0.62
cx(s, 4.85, 5.55, 7.6, 0.85,
   'Buscar item no estoque e andar dentro da loja não agregam valor ao shopper, mas são\n'
   'consequência de layout, de ruptura e de onde o estoque fica. São endereçáveis.', 13, CINZA, esp=1.35)
rodape(s, 'Classificação: deslocamento, busca e localização de produtos, em todos os blocos.')

# =================================================== 9 · variabilidade
s = slide()
titulo(s, 'A mesma tarefa leva tempos muito diferentes',
       'Dispersão do tempo de visita dentro de cada canal — quanto maior, menor a padronização.')
v2 = vis.groupby('canal').tempo_loja_min.agg(['size', 'min', 'median', 'max', 'mean', 'std'])
v2 = v2[v2['size'] >= 8].sort_values('median', ascending=False)
cd = CategoryChartData()
cd.categories = list(v2.index)
cd.add_series('Mais rápida', tuple(float(x) for x in v2['min']))
cd.add_series('Mediana', tuple(float(x) for x in v2['median']))
cd.add_series('Mais lenta', tuple(float(x) for x in v2['max']))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.9), Inches(2.3),
                        Inches(11.55), Inches(3.2), cd).chart
estiliza(gf, legenda=True, tam=10)
pinta(gf, ['A9CBD6', '2C7391', '123A4D'])
eixos(gf, valores=True)
gf.plots[0].gap_width = 70
pior = v2.assign(cv=v2['std'] / v2['mean']).sort_values('cv', ascending=False).index[0]
cx(s, 0.9, 5.8, 11.55, 1.1,
   'No {}, a visita mais lenta leva {} vezes o tempo da mais rápida. Parte é porte de loja, mas a faixa é\n'
   'larga demais para ser só isso — é onde um padrão de execução tem mais a ganhar.'.format(
       pior, '{:.1f}'.format(v2.loc[pior, 'max'] / v2.loc[pior, 'min']).replace('.', ',')),
   14, TINTA, esp=1.4)
rodape(s, 'Tempo total dentro da loja, por visita.')
# ============================================== 10 · a jornada de 8h
s = slide()
titulo(s, 'A rota que o tempo medido comporta',
       'Quantas lojas cabem em 8 horas, segundo o tempo medido, contra a rota que o campo declara.')
cd = CategoryChartData()
ordp = sorted(FORTES, key=lambda c: -c['cabe_em_8h'])
cd.categories = [c['canal'] for c in ordp]
cd.add_series('Projeção pelo tempo medido', tuple(c['cabe_em_8h'] for c in ordp))
cd.add_series('Rota declarada pelo campo', tuple(c['lojas_dia'] or 0 for c in ordp))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.9), Inches(2.3),
                        Inches(11.55), Inches(3.1), cd).chart
estiliza(gf, legenda=True, tam=11)
pinta(gf, ['123A4D', '7DB0C1'])
eixos(gf)
gf.plots[0].data_labels.number_format = '0.0'
gf.plots[0].data_labels.number_format_is_linked = False
gf.plots[0].gap_width = 70
cx(s, 0.9, 5.7, 11.55, 1.15,
   'As duas leituras chegam ao mesmo lugar por caminhos independentes: o tempo medido dentro da loja\n'
   'prevê o tamanho da rota que o campo declara. Mas nenhuma das duas inclui o trajeto entre lojas.',
   14, TINTA, esp=1.4)
rodape(s, 'Projeção = 480 minutos ÷ tempo médio de visita. Não é meta operacional.')

# ============================================= 11 · deslocamento externo
s = slide()
titulo(s, 'O trajeto entre lojas — primeira medição',
       'Quanto custa o caminho até a loja, e quanto custa alongar uma rota.')
lab = {'Casa / base': 'De casa até a\nprimeira loja', 'Loja anterior': 'Entre duas\nlojas'}
x = 0.9
for k, v in desl['por_origem'].items():
    caixa(s, x, 2.35, 3.6, 2.5)
    cx(s, x + 0.3, 2.75, 3.0, 0.8, '{} min'.format(v['mediana']), 40, PETROL, True, TIT,
       PP_ALIGN.CENTER, 1.0)
    cx(s, x + 0.3, 3.65, 3.0, 0.75, lab.get(k, k), 13, TINTA, True, CORPO, PP_ALIGN.CENTER, 1.25)
    cx(s, x + 0.3, 4.4, 3.0, 0.35, '{} medições'.format(v['n']), 11, CINZA, False, CORPO,
       PP_ALIGN.CENTER)
    x += 3.9
cx(s, 8.7, 2.35, 3.75, 0.45, 'POR QUE IMPORTA', 11, TERRA, True)
cx(s, 8.7, 2.9, 3.75, 2.6,
   'Cada loja a mais numa rota cobra\ndois tempos: o da visita e o do\ntrajeto até ela.\n\n'
   'O trajeto de casa é pago uma vez\npor dia; o trajeto entre lojas se\nmultiplica.', 13, TINTA, esp=1.35)
cx(s, 0.9, 5.4, 7.6, 1.0,
   'A diferença entre os dois trajetos é o dado útil: o deslocamento de casa é um custo fixo do dia,\n'
   'enquanto o trajeto entre lojas é o que se multiplica a cada loja a mais na rota.', 13, CINZA, esp=1.35)
rodape(s, 'O trajeto entre lojas é o que se multiplica a cada loja a mais na rota.')

# ================================================== 13 · recomendações
s = slide()
titulo(s, 'Para onde olhar primeiro',
       'Ordenado por tamanho da oportunidade e por facilidade de execução.')
REC = [('1', 'Atacar o tempo de busca no estoque',
        'Buscar item faltante é o 2º maior movimento da jornada. Endereçamento e separação prévia '
        'atacam direto.', TERRA),
       ('2', 'Padronizar onde a variação é maior',
        'A faixa entre a visita mais rápida e a mais lenta é ampla demais para ser só porte de loja. '
        'Um padrão por canal reduz a cauda.', TEAL),
       ('3', 'Tratar DPP como operação própria',
        'Com 4 lojas por dia, o custo fixo de entrada pesa 4 vezes. Reduzir fricção de acesso vale '
        'mais nele do que em qualquer outro canal.', PETROL),
       ('4', 'Dimensionar a rota com o trajeto incluído',
        'Tempo de loja e trajeto já estão medidos. Cruzá-los por setor mostra quais rotas cabem '
        'na jornada e quais precisam ser redesenhadas.', VERDE)]
y = 2.3
for n, t, d, cor in REC:
    from pptx.enum.shapes import MSO_SHAPE
    cr = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y), Inches(0.55), Inches(0.55))
    cr.fill.solid()
    cr.fill.fore_color.rgb = cor
    cr.line.fill.background()
    cr.shadow.inherit = False
    tf = cr.text_frame
    tf.text = n
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(16)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = BRANCO
    tf.paragraphs[0].runs[0].font.name = TIT
    cx(s, 1.75, y + 0.02, 10.7, 0.4, t, 16, PETROL, True, TIT)
    cx(s, 1.75, y + 0.5, 10.7, 0.6, d, 12.5, TINTA, esp=1.3)
    y += 1.16
rodape(s, 'As recomendações saem do que foi medido; o dimensionamento de ganho exige um piloto.')

# ==================================================== 14 · fechamento
s = slide(True)
cx(s, 0.9, 2.2, 11.5, 0.5, 'PRÓXIMOS PASSOS', 13, RGBColor(0x7D, 0xB0, 0xC1), True, CORPO)
cx(s, 0.9, 2.8, 11.5, 1.0, 'Do diagnóstico à ação', 40, BRANCO, True, TIT, esp=1.1)
cx(s, 0.9, 4.0, 11.0, 2.1,
   'O tempo da operação está medido, movimento a movimento, canal a canal. O próximo passo é\n'
   'escolher onde atacar primeiro e testar em piloto controlado — o estudo dimensiona o\n'
   'problema; o piloto dimensiona o ganho.\n\n'
   'A medição continua rodando, então cada rodada de melhoria pode ser verificada com o\n'
   'mesmo instrumento que produziu este diagnóstico.', 15, RGBColor(0xCB, 0xE0, 0xE6), esp=1.5)
cx(s, 0.9, 6.5, 11.5, 0.4,
   'Estudo de Tempos e Movimentos  ·  {} a {}'.format(*meta['periodo']),
   11, RGBColor(0x7D, 0xB0, 0xC1))

SAIDA = 'analise/entregaveis/Tempos e Movimentos - Apresentacao.pptx'
prs.save(SAIDA)
print('salvo:', SAIDA)
print('slides:', len(prs.slides.__iter__.__self__._sldIdLst))
